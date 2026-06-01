"""
破冰协议 · AI 引导 demo 后端

启动方式：
  uv run icebreaker-demo

环境变量（或 .env 文件，可选）：
  DEEPSEEK_API_KEY=sk-xxx
  DEEPSEEK_BASE_URL=https://api.deepseek.com
  DEEPSEEK_MODEL=deepseek-chat

没有 API key 时，后端会使用内置规则引擎，保证 demo 仍然能跑通。
"""

from __future__ import annotations

import json
import os
import re
import base64
from io import BytesIO
from pathlib import Path
from typing import Literal

from dotenv import load_dotenv
from fastapi import FastAPI, Request
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from openai import OpenAI
from pydantic import BaseModel, Field
from starlette.middleware.base import BaseHTTPMiddleware

load_dotenv()


API_KEY = os.getenv("DEEPSEEK_API_KEY", "")
BASE_URL = os.getenv("DEEPSEEK_BASE_URL", "https://api.deepseek.com")
MODEL = os.getenv("DEEPSEEK_MODEL", "deepseek-chat")

client = OpenAI(api_key=API_KEY, base_url=BASE_URL) if API_KEY else None
HOST = "0.0.0.0"
PORT = 8000


Screen = Literal["contract", "roadmap", "step", "done", "message", "thinking_budget"]


class Step(BaseModel):
    title: str
    instruction: str
    output: str
    minutes: int = Field(ge=1, le=15)


class Attachment(BaseModel):
    name: str
    type: str | None = None
    size: int | None = None
    text: str | None = None


class ParseAttachmentRequest(BaseModel):
    name: str
    type: str | None = None
    data_base64: str


class ParseAttachmentResponse(BaseModel):
    name: str
    type: str | None = None
    size: int
    text: str = ""
    parsed: bool = False
    error: str | None = None


class ChatRequest(BaseModel):
    message: str
    history: list[dict[str, str]] = Field(default_factory=list)
    phase: str | None = None
    task: str | None = None
    steps: list[Step] = Field(default_factory=list)
    current_step: int = 0
    outputs: list[str] = Field(default_factory=list)
    attachments: list[Attachment] = Field(default_factory=list)


class ChatResponse(BaseModel):
    reply: str
    screen: Screen = "message"
    task: str | None = None
    steps: list[Step] = Field(default_factory=list)
    current_step: int | None = None
    thinking_budget_seconds: int | None = None


class UTF8JSONResponse(JSONResponse):
    media_type = "application/json; charset=utf-8"


CONTRACT_REPLY = (
    "先签破冰契约：这次目标不是完美，而是做出一个能改的雏形。"
    "每一步都要有看得见的产出；允许不满意，但不允许空白。"
)

SHAME_REFRAMES = {
    "怕做不好": "第一版不需要完美，但必须存在。先拿到一个能改的对象。",
    "别人会怎么看": "当前没有别人。初稿是给你自己改的，不是给别人审判的。",
    "我不行": "不是你不行，是你还没有开始。行不行得做了才知道。",
    "想清楚再做": "想清楚是无终止条件。先写出一个版本，再用版本继续思考。",
    "做得太烂": "初稿的价值不在于好不好，在于它给了你一个可以修改的对象。",
}


def normalize_text(value: str) -> str:
    return re.sub(r"\s+", " ", value or "").strip()


def is_agreement(message: str) -> bool:
    text = normalize_text(message).lower()
    return any(token in text for token in ("同意", "开始", "拆解", "继续", "ok", "好", "go"))


def is_done_request(message: str) -> bool:
    text = normalize_text(message)
    return any(token in text for token in ("所有步骤", "全部完成", "都完成", "拼装成型", "雏形完成"))


def is_thinking_delay(message: str) -> bool:
    text = normalize_text(message)
    return any(token in text for token in ("没想清楚", "再想想", "还没准备", "不知道怎么开始", "准备好再"))


def is_stuck(message: str) -> bool:
    text = normalize_text(message)
    return any(token in text for token in ("卡住", "不会", "太难", "不知道怎么写", "下不了手", "没有思路"))


def extract_task(req: ChatRequest) -> str:
    if req.task:
        return req.task

    ignored = ("同意", "开始", "拆解", "下一步", "做完", "所有步骤", "拼装")
    for item in req.history:
        if item.get("role") != "user":
            continue
        content = normalize_text(item.get("content", ""))
        if content and not any(token in content for token in ignored):
            return content

    return normalize_text(req.message) or "做出一个最小雏形"


def detect_category(task: str) -> str:
    text = normalize_text(task)
    if any(token in text for token in ("认识", "开口", "联系", "消息", "社交", "同行")):
        return "social"
    if any(token in text for token in ("简历", "求职", "面试", "jd", "岗位", "投递")):
        return "job"
    if any(token in text for token in ("代码", "编程", "python", "项目", "demo", "网站", "爬虫", "app")):
        return "code"
    if any(token in text for token in ("文章", "博客", "小红书", "脚本", "笔记", "写")):
        return "content"
    return "general"


def build_rule_steps(task: str) -> list[Step]:
    category = detect_category(task)

    if "不知道" in task and category in {"code", "content", "general"}:
        return [
            Step(
                title="回答三个约束",
                instruction="写下三个二选一答案：时间是一整天还是2小时；能量是高还是低；目的是学习还是产出。",
                output="三个约束答案",
                minutes=5,
            ),
            Step(
                title="选一个安全方向",
                instruction="从安全、探索、冒险里选一个。选不出来就默认安全方向，别在选项里继续打转。",
                output="一个被选中的方向",
                minutes=5,
            ),
            Step(
                title="做最小可见物",
                instruction="围绕这个方向做一个最小可见物：一段文字、一个文件、一个页面或一条可发送消息。",
                output="最小可见物",
                minutes=12,
            ),
        ]

    if category == "social":
        return [
            Step(
                title="锁定对象和目的",
                instruction="写下你要联系谁，以及你只想完成的一个目的。不要评价自己配不配。",
                output="对象 + 一个联系目的",
                minutes=5,
            ),
            Step(
                title="生成可复制消息",
                instruction="用这个模板改成一条能直接复制的消息：hi，看到你在做[方向]，我也在关注相关内容，方便加个微信交流一下吗？",
                output="一条完整消息",
                minutes=5,
            ),
            Step(
                title="发送或保存发送稿",
                instruction="复制这条消息。能发就发；今天实在发不出去，就把发送稿保存到一个固定位置。",
                output="已发送状态或保存好的发送稿",
                minutes=2,
            ),
        ]

    if category == "job":
        return [
            Step(
                title="摘出一个岗位",
                instruction="只选一个岗位，把岗位名和链接写下来。不要比较十个岗位。",
                output="一个目标岗位",
                minutes=5,
            ),
            Step(
                title="圈三个关键词",
                instruction="从 JD 里圈出三个你能回应的关键词。不是优势分析，只看你做过的具体事。",
                output="三个 JD 关键词",
                minutes=8,
            ),
            Step(
                title="改一句简历",
                instruction="把简历里一句经历改成包含这三个关键词的版本。只改一句。",
                output="一条可投递的经历描述",
                minutes=8,
            ),
            Step(
                title="完成一次投递动作",
                instruction="把这份简历用于一次投递，或把投递材料打包到同一个文件夹。",
                output="一次投递或投递包",
                minutes=5,
            ),
        ]

    if category == "code":
        return [
            Step(
                title="写一句 demo 目标",
                instruction="用一句话写清楚这个 demo 第一版要让用户看到什么。只允许一个核心动作。",
                output="一句 demo 目标",
                minutes=5,
            ),
            Step(
                title="做出最小骨架",
                instruction="创建最少文件，让页面、脚本或接口能打开。里面可以只有一个真实按钮或一条真实输出。",
                output="可打开的最小骨架",
                minutes=10,
            ),
            Step(
                title="跑通核心动作",
                instruction="让那个唯一核心动作跑通一次，并记录看到的结果。不要加第二个功能。",
                output="一次跑通记录",
                minutes=12,
            ),
            Step(
                title="写下下一处改进",
                instruction="只写一个下一轮最值得改的点。写完就停。",
                output="一个明确改进点",
                minutes=5,
            ),
        ]

    if category == "content":
        return [
            Step(
                title="写一句核心观点",
                instruction="写一句话说明你想表达什么。句子难看没关系，但必须是一句话。",
                output="一句核心观点",
                minutes=5,
            ),
            Step(
                title="写三个标题",
                instruction="围绕这句话写 3 个标题。不要判断好坏，先列出来。",
                output="三个标题",
                minutes=5,
            ),
            Step(
                title="搭三段结构",
                instruction="写出开头、主体、结尾各一句。每段只写一句。",
                output="三段结构",
                minutes=8,
            ),
            Step(
                title="扩成初稿",
                instruction="把三段各扩写成一小段，凑成完整初稿。",
                output="完整初稿",
                minutes=12,
            ),
        ]

    return [
        Step(
            title="定义雏形",
            instruction="用一句话写下这件事完成到什么程度就算有雏形。",
            output="一句雏形定义",
            minutes=5,
        ),
        Step(
            title="列出三个部件",
            instruction="写出这个雏形需要的三个最小部件。不要超过三个。",
            output="三个最小部件",
            minutes=8,
        ),
        Step(
            title="做第一个部件",
            instruction="只做第一个部件，让它变成看得见的东西。",
            output="第一个可见部件",
            minutes=12,
        ),
        Step(
            title="拼成初稿记录",
            instruction="把已经有的东西整理成一段记录，并写下下一步只改哪里。",
            output="初稿记录 + 一个改进点",
            minutes=10,
        ),
    ]


def safe_history(history: list[dict[str, str]]) -> list[dict[str, str]]:
    cleaned: list[dict[str, str]] = []
    for item in history[-12:]:
        role = item.get("role")
        content = normalize_text(item.get("content", ""))
        if role in {"user", "assistant"} and content:
            cleaned.append({"role": role, "content": content[:2000]})
    return cleaned


def attachment_context(attachments: list[Attachment]) -> str:
    if not attachments:
        return ""

    lines = ["", "[附件摘要]"]
    for item in attachments[:8]:
        name = normalize_text(item.name)[:120] or "未命名附件"
        file_type = normalize_text(item.type or "unknown")
        size = f"{item.size} bytes" if item.size else "unknown size"
        text = normalize_text(item.text or "")
        lines.append(f"- {name} ({file_type}, {size})")
        if text:
            lines.append(text[:1200])
    return "\n".join(lines)[:5000]


def try_ai_plan(task: str, attachments: list[Attachment] | None = None) -> list[Step] | None:
    if client is None:
        return None

    attached_context = attachment_context(attachments or [])
    prompt = f"""
你是破冰协议的任务拆解器。用户任务：{task}

请只返回 JSON，不要 Markdown：
{{
  "steps": [
    {{
      "title": "短标题",
      "instruction": "具体执行指令，不能是想一想，必须让用户产出东西",
      "output": "做完后可见的产出物",
      "minutes": 5
    }}
  ]
}}

规则：
- 3 到 5 步。
- 每步 1 到 15 分钟。禁止生成超过 15 分钟的步骤，任务复杂就拆成多个原子步骤。
- 每步必须有可见产出。
- 第一版目标是可修改的雏形，不是完美成品。
"""
    if attached_context:
        prompt += attached_context
    try:
        response = client.chat.completions.create(
            model=MODEL,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=700,
            temperature=0.2,
        )
        raw = response.choices[0].message.content or ""
        match = re.search(r"\{.*\}", raw, flags=re.S)
        payload = json.loads(match.group(0) if match else raw)
        steps = [Step(**item) for item in payload.get("steps", [])]
        return steps if 3 <= len(steps) <= 6 else None
    except Exception:
        return None


def try_ai_reply(req: ChatRequest) -> str | None:
    if client is None:
        return None

    system_prompt = (
        '你是 [Protocol]，破冰协议的 AI 助手。用户正在一个分步工作流中完成任务。\n'
        '你会收到完整的上下文：用户的总任务、已完成步骤的产出、当前步骤要求。\n\n'
        '规则：\n'
        '- 你的名字是 [Protocol]，以此身份回复。\n'
        '- 不要安慰、不要讲道理、不要评价用户。\n'
        '- 直接产出内容：如果用户要标题就给标题，要代码就给代码，要消息就给消息。\n'
        '- 给出可以直接使用的具体内容，而不是「你可以试试...」这种空泛建议。\n'
        '- 如果用户没指定方向，基于上下文主动给出最合理的版本。\n'
        '- 回复简洁，但可以包含多个选项供用户挑选。\n'
        '- 排版清晰：用换行分隔不同内容块，选项用编号列表。'
    )
    messages = [{"role": "system", "content": system_prompt}]
    messages.extend(safe_history(req.history))
    messages.append({"role": "user", "content": (req.message[:2000] + attachment_context(req.attachments))[:7000]})

    try:
        response = client.chat.completions.create(
            model=MODEL,
            messages=messages,
            max_tokens=500,
            temperature=0.5,
        )
        return (response.choices[0].message.content or "").strip()
    except Exception:
        return None


def build_roadmap(req: ChatRequest) -> ChatResponse:
    task = extract_task(req)
    steps = try_ai_plan(task, req.attachments) or build_rule_steps(task)
    reply = f"任务先收窄：{task}。我把它拆成 {len(steps)} 步，每步都有一个可见产出。我们只从第 1 步开始。"
    return ChatResponse(reply=reply, screen="roadmap", task=task, steps=steps, current_step=0)


def build_step_help(req: ChatRequest) -> ChatResponse:
    current_index = max(0, min(req.current_step, len(req.steps) - 1)) if req.steps else 0
    step = req.steps[current_index] if req.steps else None

    if is_thinking_delay(req.message):
        return ChatResponse(
            reply="你的思考预算是 5 分钟。时间到以后，不管满意不满意，先写出一个版本。",
            screen="thinking_budget",
            task=req.task,
            steps=req.steps,
            current_step=current_index,
            thinking_budget_seconds=300,
        )

    ai_reply = try_ai_reply(req)
    if ai_reply:
        return ChatResponse(reply=ai_reply, screen="message", task=req.task, steps=req.steps, current_step=current_index)

    if step:
        reply = f"别跳走。当前只做这一格：{step.instruction} 做完你只需要交出：{step.output}。"
    else:
        reply = "先写出一个最小版本。不要解释为什么还没准备好，直接留下一个可见产出。"
    return ChatResponse(reply=reply, screen="message", task=req.task, steps=req.steps, current_step=current_index)


def build_done(req: ChatRequest) -> ChatResponse:
    visible_outputs = [normalize_text(item) for item in req.outputs if normalize_text(item)]
    if visible_outputs:
        reply = f"雏形已经成立：你产出了 {len(visible_outputs)} 个可见块。下一轮只改一个地方，别开新战场。"
    else:
        reply = "雏形还很薄，但流程已经跑通。下一轮先补一个真实产出，再谈优化。"
    return ChatResponse(reply=reply, screen="done", task=req.task, steps=req.steps, current_step=len(req.steps))


class NoCacheMiddleware(BaseHTTPMiddleware):
    """HTML/SW 禁用缓存；CSS/JS 短缓存 + must-revalidate。"""

    async def dispatch(self, request: Request, call_next):
        response = await call_next(request)
        ct = response.headers.get("content-type", "")
        path = request.url.path
        # HTML + Service Worker 文件必须永不缓存
        if "text/html" in ct or path.endswith("/sw.js") or path.endswith("/registerSW.js"):
            response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate, private"
            response.headers["Pragma"] = "no-cache"
            response.headers["Expires"] = "0"
        elif "javascript" in ct or "css" in ct:
            response.headers["Cache-Control"] = "max-age=300, must-revalidate"
        return response


app = FastAPI(title="破冰协议 API", default_response_class=UTF8JSONResponse)
app.add_middleware(NoCacheMiddleware)


@app.post("/api/chat", response_model=ChatResponse)
async def chat(req: ChatRequest):
    message = normalize_text(req.message)

    for trigger, reply in SHAME_REFRAMES.items():
        if trigger in message:
            return ChatResponse(reply=reply, screen="message", task=req.task, steps=req.steps, current_step=req.current_step)

    if is_done_request(message):
        return build_done(req)

    if req.phase == "step" or is_stuck(message) or is_thinking_delay(message):
        return build_step_help(req)

    if is_agreement(message):
        return build_roadmap(req)

    task = extract_task(req)
    return ChatResponse(reply=CONTRACT_REPLY, screen="contract", task=task)


def _fallback_reply(req: ChatRequest) -> str:
    """AI 不可用时，用规则引擎生成回复。"""
    message = normalize_text(req.message)
    for trigger, reply in SHAME_REFRAMES.items():
        if trigger in message:
            return reply
    current_index = max(0, min(req.current_step, len(req.steps) - 1)) if req.steps else 0
    step = req.steps[current_index] if req.steps else None
    if step:
        return f"别跳走。当前只做这一格：{step.instruction} 做完你只需要交出：{step.output}。"
    return "先写出一个最小版本。不要解释为什么还没准备好，直接留下一个可见产出。"


def stream_ai_reply(req: ChatRequest):
    """SSE 流式生成：逐 token 输出 AI 回复，最后发送 [DONE] 标记。"""
    if client is None:
        fallback = _fallback_reply(req)
        yield f"data: {json.dumps({'text': fallback})}\n\n"
        yield "data: [DONE]\n\n"
        return

    system_prompt = (
        '你是 [Protocol]，破冰协议的 AI 助手。用户正在一个分步工作流中完成任务。\n'
        '你会收到完整的上下文：用户的总任务、已完成步骤的产出、当前步骤要求。\n\n'
        '规则：\n'
        '- 你的名字是 [Protocol]，以此身份回复。\n'
        '- 不要安慰、不要讲道理、不要评价用户。\n'
        '- 直接产出内容：如果用户要标题就给标题，要代码就给代码，要消息就给消息。\n'
        '- 给出可以直接使用的具体内容，而不是「你可以试试...」这种空泛建议。\n'
        '- 如果用户没指定方向，基于上下文主动给出最合理的版本。\n'
        '- 回复简洁，但可以包含多个选项供用户挑选。\n'
        '- 排版清晰：用换行分隔不同内容块，选项用编号列表。'
    )
    messages = [{"role": "system", "content": system_prompt}]
    messages.extend(safe_history(req.history))
    messages.append({"role": "user", "content": (req.message[:2000] + attachment_context(req.attachments))[:7000]})

    try:
        stream = client.chat.completions.create(
            model=MODEL,
            messages=messages,
            max_tokens=500,
            temperature=0.5,
            stream=True,
        )
        for chunk in stream:
            delta = chunk.choices[0].delta
            if delta.content:
                yield f"data: {json.dumps({'text': delta.content})}\n\n"
        yield "data: [DONE]\n\n"
    except Exception:
        fallback = _fallback_reply(req)
        yield f"data: {json.dumps({'text': fallback})}\n\n"
        yield "data: [DONE]\n\n"


@app.post("/api/chat/stream")
async def chat_stream(req: ChatRequest):
    return StreamingResponse(
        stream_ai_reply(req),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


# ==================== 摘要归档 ====================

def parse_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(BytesIO(data))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def parse_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(BytesIO(data))
    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def parse_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets[:10]:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(max_row=200, values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                parts.append("\t".join(values).rstrip())
    return "\n".join(parts)


def parse_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(BytesIO(data))
    parts: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            parts.append(f"# Slide {index}\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


def parse_office_attachment(name: str, data: bytes) -> tuple[str, bool, str | None]:
    suffix = Path(name).suffix.lower()
    try:
        if suffix == ".pdf":
            return parse_pdf(data), True, None
        if suffix == ".docx":
            return parse_docx(data), True, None
        if suffix == ".xlsx":
            return parse_xlsx(data), True, None
        if suffix == ".pptx":
            return parse_pptx(data), True, None
        return "", False, "unsupported_type"
    except Exception as exc:
        return "", False, str(exc)[:200]


@app.post("/api/attachments/parse", response_model=ParseAttachmentResponse)
async def parse_attachment(req: ParseAttachmentRequest):
    try:
        data = base64.b64decode(req.data_base64, validate=True)
    except Exception:
        return ParseAttachmentResponse(name=req.name, type=req.type, size=0, error="invalid_base64")

    if len(data) > 15 * 1024 * 1024:
        return ParseAttachmentResponse(name=req.name, type=req.type, size=len(data), error="file_too_large")

    text, parsed, error = parse_office_attachment(req.name, data)
    text = normalize_text(text)[:60000]
    return ParseAttachmentResponse(
        name=req.name,
        type=req.type,
        size=len(data),
        text=text,
        parsed=parsed and bool(text),
        error=error if not text else None,
    )


class SummarizeRequest(BaseModel):
    step_title: str
    user_content: str


class SummarizeResponse(BaseModel):
    summary: str


SUMMARIZE_PROMPT = """你是破冰协议的归档引擎。

用户刚完成步骤「{step_title}」，提交了以下内容：
{user_content}

请将其压缩成一句话（不超过20字），用于破冰战报的归档展示。
不要评价好坏，不要给建议，只提炼核心信息。
如果内容很乱，就从中找出最关键的那个词或那件事。

只返回那一句话，不要任何其他内容。"""


@app.post("/api/summarize", response_model=SummarizeResponse)
async def summarize(req: SummarizeRequest):
    if client is None:
        return SummarizeResponse(summary=req.user_content.strip()[:20])

    prompt = SUMMARIZE_PROMPT.format(
        step_title=req.step_title,
        user_content=req.user_content[:1000],
    )
    try:
        response = client.chat.completions.create(
            model=MODEL,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=50,
            temperature=0.2,
        )
        summary = (response.choices[0].message.content or "").strip()
        return SummarizeResponse(summary=summary)
    except Exception:
        return SummarizeResponse(summary=req.user_content.strip()[:20])


REVIEW_DIMENSIONS = [
    {"key": "completion", "name": "完成度", "desc": "是否交付了每一步要求的可见产出"},
    {"key": "clarity", "name": "清晰度", "desc": "别人能否快速理解产出在说什么、要做什么"},
    {"key": "usefulness", "name": "可用性", "desc": "当前版本是否已经能被继续使用、修改或展示"},
    {"key": "audience_fit", "name": "受众匹配", "desc": "是否命中目标受众或使用场景"},
    {"key": "next_action", "name": "下一步明确度", "desc": "是否清楚下一刀应该改哪里"},
]

REVIEW_PROMPT = """你是破冰协议的产出质量评审员。
评价用户的产出，不评价用户。

只返回 JSON，不要 Markdown，不要 JSON 以外的文字。
JSON 中所有自然语言值必须用简体中文。

固定维度：
{dimensions}

JSON 结构：
{{
  "verdict": "一句话结论",
  "summary": "2-3 句话描述当前质量水平",
  "strengths": ["最多 3 个值得保留的优点"],
  "issues": ["最多 4 个产出的具体问题"],
  "priority_fix": "40 字以内的一个可执行修改指令",
  "dimensions": [
    {{"key":"completion","score":1,"comment":"原因"}},
    {{"key":"clarity","score":1,"comment":"原因"}},
    {{"key":"usefulness","score":1,"comment":"原因"}},
    {{"key":"audience_fit","score":1,"comment":"原因"}},
    {{"key":"next_action","score":1,"comment":"原因"}}
  ]
}}

规则：
- 不要自我评价，只评价产出。
- 严格且可执行。
- 不要安慰或空洞赞美。
- 每个问题必须指向产出本身。
- priority_fix 必须是唯一的最佳下一步修改。

待评价内容：
{payload}"""


class ReviewStep(BaseModel):
    index: int
    title: str
    instruction: str = ""
    expected_output: str = ""
    user_output: str = ""


class ReviewRequest(BaseModel):
    task: str = ""
    steps: list[ReviewStep] = Field(default_factory=list)
    session_log: str | None = None
    output_mode: str | None = None
    protocol_strength: str | None = None


class ReviewDimension(BaseModel):
    key: str
    name: str = ""
    desc: str = ""
    score: int = Field(ge=1, le=5)
    comment: str = ""


class ReviewResponse(BaseModel):
    total: int = Field(ge=0)
    max: int = Field(default=25)
    verdict: str = ""
    summary: str = ""
    strengths: list[str] = Field(default_factory=list)
    issues: list[str] = Field(default_factory=list)
    priority_fix: str = ""
    dimensions: list[ReviewDimension] = Field(default_factory=list)
    mode: str = "ai"
    error: str | None = None


def _fallback_review(payload: ReviewRequest) -> dict:
    steps = payload.steps or []
    filled = sum(1 for s in steps if len(normalize_text(s.user_output)) >= 5)
    total_steps = max(len(steps), 1)
    base = max(2, min(4, round((filled / total_steps) * 5)))
    dimensions = [
        {**dim, "score": base, "comment": "已有可见产出，但还需要更具体的验收标准和表达打磨。"
         if filled == total_steps else "部分步骤缺少足够清晰的产出，先补齐空白再谈优化。"}
        for dim in REVIEW_DIMENSIONS
    ]
    if filled < total_steps:
        dimensions[0]["score"] = 2
        dimensions[0]["comment"] = "部分步骤还缺少足够具体的产出。"
    return {
        "total": sum(d["score"] for d in dimensions),
        "max": len(dimensions) * 5,
        "verdict": "能继续打磨" if filled == total_steps else "需要补齐产出",
        "summary": "AI评价暂时不可用，已根据完成步骤做本地兜底判断。",
        "strengths": ["已经完成了破冰流程，至少留下了可修改的版本。"],
        "issues": ["需要补充更明确的边界、验收标准和面向受众的表达。"],
        "priority_fix": "先补齐最薄弱的一步：让它有一个别人能看懂的具体产出。",
        "dimensions": dimensions,
    }


def _ai_review(payload: ReviewRequest) -> dict:
    if client is None:
        return {"error": "api_key_missing"}

    prompt = REVIEW_PROMPT.format(
        dimensions=json.dumps(REVIEW_DIMENSIONS, ensure_ascii=False),
        payload=json.dumps(payload.model_dump(), ensure_ascii=False)[:9000],
    )
    try:
        response = client.chat.completions.create(
            model=MODEL,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1200,
            temperature=0.2,
            timeout=30.0,
        )
        raw = (response.choices[0].message.content or "").strip()
        match = re.search(r"\{.*\}", raw, flags=re.S)
        if not match:
            return {"error": "invalid_response", "detail": "AI 返回内容无法解析为 JSON"}
        return json.loads(match.group(0))
    except Exception as exc:
        import logging
        logging.getLogger("icebreaker").warning("AI review failed: %s", exc)
        error_type = "invalid_response"
        msg = str(exc).lower()
        name = type(exc).__name__.lower()
        if "timeout" in name or "timeout" in msg:
            error_type = "timeout"
        elif "auth" in name or "401" in msg or "403" in msg:
            error_type = "api_key_invalid"
        elif "connect" in name or "connect" in msg:
            error_type = "connection_error"
        return {"error": error_type, "detail": str(exc)[:200]}


@app.post("/api/review", response_model=ReviewResponse)
async def review(req: ReviewRequest):
    result = _ai_review(req)

    if result and "error" not in result:
        # 钳制 AI 返回的分数到 1-5 范围
        for dim in result.get("dimensions", []):
            if "score" in dim:
                dim["score"] = max(1, min(5, int(dim["score"])))
        # 补全 AI 可能省略的字段
        dims = result.get("dimensions", [])
        result.setdefault("total", sum(d.get("score", 0) for d in dims))
        result.setdefault("max", len(dims) * 5 if dims else 25)
        result.setdefault("verdict", "")
        result.setdefault("summary", "")
        result.setdefault("strengths", [])
        result.setdefault("issues", [])
        result.setdefault("priority_fix", "")
        result.setdefault("mode", "ai")
        return result

    error_type = result.get("error", "unknown") if result else "unknown"
    fallback = _fallback_review(req)
    fallback["mode"] = "local"
    fallback["error"] = error_type
    return fallback


@app.get("/api/key-status")
async def key_status():
    """检查 AI API Key 配置状态和连通性。"""
    if not client:
        return {"configured": False}
    try:
        client.chat.completions.create(
            model=MODEL,
            messages=[{"role": "user", "content": "hi"}],
            max_tokens=5,
            timeout=10.0,
        )
        return {"configured": True, "valid": True}
    except Exception:
        return {"configured": True, "valid": False}


demo_dir = Path(__file__).parent / "demo"


# 优先使用预构建的 dist/ 目录（Vite 生产构建产物）
# 若不存在则回退到 demo/ 源码目录（Vite dev server 场景）
dist_dir = demo_dir / "dist"
if dist_dir.is_dir():
    # API 路由已在此前注册，挂载到根路径兜底静态文件和 SPA fallback
    app.mount("/", StaticFiles(directory=dist_dir, html=True), name="static")


def run_demo() -> None:
    import uvicorn

    print("\n  破冰协议 · AI 引导 demo")
    print(f"  http://localhost:{PORT}\n")
    uvicorn.run(app, host=HOST, port=PORT)


if __name__ == "__main__":
    run_demo()
